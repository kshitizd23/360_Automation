import os
import io
import pandas as pd
from flask import Flask, render_template, request, send_file, jsonify, abort
from pptx import Presentation
from pptx.chart.data import CategoryChartData

# from pptx.enum.chart import XL_CHART_TYPE # Not needed for replacing data
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor          # For specifying colors
from pptx.enum.chart import XL_MARKER_STYLE # For marker shapes

import logging
import openpyxl # For Excel override

# --- Configuration ---
DATA_FILE = 'data.csv'
TEMPLATE_FILE = 'template2.pptx' # Your pre-built report
OVERRIDE_EXCEL_FILE = 'Chart_Overrides.xlsx' # For overrides
OUTPUT_FILENAME = 'Updated_Report.pptx'

# --- Naming Conventions (Must Match PowerPoint Selection Pane) ---
SLIDE_TITLE_SHAPE_PREFIX = "TITLE_"
TA_NAME_TEXTBOX_NAME = "TA_NAME"
MAC_CHART_SUFFIX = "_MAC"
NPS_CHART_SUFFIX = "_BRAND_NPS"
EQUITY_CHART_SUFFIX = "_BRAND_EQUITY"
# New Suffix/Convention for MVP Charts
MVP_CHART_IDENTIFIER = "_MVP_" # e.g., TA + MVP_CHART_IDENTIFIER + Country_ID

# --- Logging Setup ---
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# --- Flask App Initialization ---
app = Flask(__name__)
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0 # Prevent caching during development

# --- Global Data Store ---
df_global = None
unique_waves = []
unique_tas = []
unique_countries = []

# --- Helper Functions ---

def clean_net_score(score):
    """Removes '%' and converts to float (decimal), returns NaN on error."""
    if isinstance(score, (int, float)):
        return float(score)
    if isinstance(score, str):
        try:
            return float(score.replace('%', '').strip()) / 100.0
        except ValueError:
            return pd.NA
    return pd.NA

def parse_dynamic_form_data(form_dict):
    """Parses brand/company selections from form data."""
    ta_selections = {}
    selected_tas = form_dict.getlist('ta')
    for ta in selected_tas:
        brand_field = f'brands_{ta}'
        company_field = f'companies_{ta}'
        ta_selections[ta] = {
            'brands': form_dict.getlist(brand_field),
            'companies': form_dict.getlist(company_field)
        }
    logging.info(f"Parsed TA selections: {ta_selections}")
    return ta_selections

def find_slide_by_title_shape(presentation, target_title_shape_name):
    """Finds a slide by searching for a shape with a specific name."""
    logging.debug(f"Searching for slide with title shape: '{target_title_shape_name}'")
    for i, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if shape.name == target_title_shape_name:
                logging.info(f"Found slide for '{target_title_shape_name}' at index {i}")
                return slide
    logging.warning(f"Slide with title shape '{target_title_shape_name}' not found in the presentation.")
    return None

def prepare_chart_data(df_pivot):
    """Creates a CategoryChartData object from a pivot table."""
    chart_data = CategoryChartData()
    if df_pivot is None or df_pivot.empty:
        logging.warning("prepare_chart_data called with empty or None pivot data. Returning empty ChartData.")
        return chart_data

    try:
        # Categories from first column
        if df_pivot.shape[1] > 0:
            chart_data.categories = df_pivot.iloc[:, 0].astype(str).tolist()
        else:
            logging.warning("Pivot table has no columns for categories.")
            return chart_data

        # Series from remaining columns
        if df_pivot.shape[1] > 1:
            for col_name in df_pivot.columns[1:]:
                # Fill NaN with 0 before converting to list for the chart
                series_values = df_pivot[col_name].fillna(0).tolist()
                chart_data.add_series(str(col_name), series_values)
        else:
            logging.warning("Pivot table has no data columns for series.")

    except Exception as e:
        logging.error(f"Error preparing chart data: {e}", exc_info=True)
        return CategoryChartData() # Return fresh empty on error

    return chart_data

def find_shape_on_slide(slide, shape_name):
    """Finds a shape with a specific name on a given slide."""
    if slide is None:
        logging.warning(f"Cannot search for shape '{shape_name}', slide object is None.")
        return None
    # Use slide.name if available (requires pptx version >= 0.6.22 approx)
    slide_identifier = getattr(slide, 'name', f"ID {getattr(slide, 'slide_id', 'N/A')}")
    logging.debug(f"Searching for shape '{shape_name}' on slide '{slide_identifier}'...")
    for shape in slide.shapes:
        if shape.name == shape_name:
            logging.debug(f"Found shape '{shape_name}' (ID: {shape.shape_id})")
            return shape
    logging.warning(f"Shape named '{shape_name}' not found on slide '{slide_identifier}'.")
    return None

# --- Helper Functions for Excel Override ---
def read_override_control_sheet(filepath):
    """Reads the first sheet of the override Excel file."""
    try:
        df_control = pd.read_excel(filepath, sheet_name=0, engine='openpyxl')
        logging.info(f"Successfully read control sheet from '{filepath}'.")
        if 'Chart_Name' not in df_control.columns or 'Update_Flag' not in df_control.columns:
            logging.error(f"Override control sheet missing required columns ('Chart_Name', 'Update_Flag').")
            return None
        return df_control
    except FileNotFoundError:
        logging.warning(f"Override Excel file '{filepath}' not found.")
        return None
    except Exception as e:
        logging.error(f"Error reading override control sheet from '{filepath}': {e}", exc_info=True)
        return None

def should_update(flag_value):
    """Checks if the update flag indicates an update is needed."""
    if flag_value is None: return False
    flag_str = str(flag_value).strip().upper()
    return flag_str == '1' or flag_str == 'TRUE' or flag_str == 'YES'

def find_shape_anywhere(presentation, shape_name):
    """Finds a shape by name across all slides."""
    logging.debug(f"Searching all slides for shape named '{shape_name}'...")
    for i, slide in enumerate(presentation.slides):
        shape = find_shape_on_slide(slide, shape_name)
        if shape:
            # Use slide.name if available for better logging
            slide_identifier = getattr(slide, 'name', f"index {i}")
            logging.info(f"Found shape '{shape_name}' on slide '{slide_identifier}'.")
            return shape
    logging.warning(f"Shape '{shape_name}' not found anywhere in the presentation.")
    return None

def read_named_range_data(filepath, range_name):
    """ Reads data from a Named Range in Excel using openpyxl. """
    logging.debug(f"Attempting to read named range '{range_name}' from '{filepath}'")
    try:
        workbook = openpyxl.load_workbook(filepath, data_only=True)
        if range_name not in workbook.defined_names:
            logging.warning(f"Named range '{range_name}' not found in '{filepath}'.")
            return None

        destination = workbook.defined_names[range_name].attr_text
        if '!' not in destination:
             logging.warning(f"Named range '{range_name}' destination '{destination}' invalid.")
             return None

        sheet_title, cells_range_str = destination.split('!', 1)
        sheet_title = sheet_title.strip("'")

        if sheet_title not in workbook.sheetnames:
            logging.warning(f"Sheet '{sheet_title}' for range '{range_name}' not found.")
            return None

        ws = workbook[sheet_title]
        if '!' in cells_range_str: cells_range_str = cells_range_str.split('!', 1)[1]
        cell_range = ws[cells_range_str]

        if not cell_range:
             logging.warning(f"Named range '{range_name}' seems empty.")
             return None

        data = []
        header = []
        first_row = True
        first_col_idx = cell_range[0][0].column

        for row_cells in cell_range:
            row_data = {}
            if first_row:
                header = [cell.value for cell in row_cells[1:]]
                first_row = False
                continue
            category = None
            for i, cell in enumerate(row_cells):
                 if cell.column == first_col_idx:
                     category = cell.value
                     row_data['Category'] = category
                 else:
                     col_idx_in_range = i - 1
                     if col_idx_in_range < len(header):
                          series_name = header[col_idx_in_range]
                          if series_name:
                            cell_value = cell.value
                            try:
                                if isinstance(cell_value, str) and '%' in cell_value:
                                     num_val = float(cell_value.replace('%','').strip()) / 100.0
                                elif cell_value is None: num_val = None
                                else: num_val = pd.to_numeric(cell_value, errors='coerce')
                            except ValueError:
                                 logging.warning(f"Value conversion error in range '{range_name}'. Setting NaN.")
                                 num_val = pd.NA
                            row_data[str(series_name)] = num_val
            if category is not None: data.append(row_data)

        if not data:
             logging.warning(f"No data rows extracted from named range '{range_name}'.")
             return None

        df_override = pd.DataFrame(data)
        ordered_columns = ['Category'] + [str(h) for h in header if h]
        existing_ordered_columns = [col for col in ordered_columns if col in df_override.columns]
        df_override = df_override[existing_ordered_columns]
        logging.info(f"Successfully read data from named range '{range_name}'. Shape: {df_override.shape}")
        return df_override

    except FileNotFoundError:
        logging.warning(f"Override Excel file '{filepath}' not found for range '{range_name}'.")
        return None
    except KeyError as e:
        logging.warning(f"Error accessing sheet/range for named range '{range_name}'. Invalid destination? Error: {e}")
        return None
    except Exception as e:
        logging.error(f"General error reading named range '{range_name}': {e}", exc_info=True)
        return None

# --- Country ID Mapping ---
def map_country_to_id(country_name):
    """Maps country name to its corresponding ID."""
    if country_name == 'G5': return 'G5'
    elif country_name == 'Italy - IT': return 'IT'
    elif country_name == 'Germany - DE': return 'DE'
    elif country_name == 'France - FR': return 'FR'
    elif country_name == 'Spain - ES': return 'ES'
    elif country_name == 'Great Britain - GB': return 'UK'
    else: return country_name # Default

# --- Data Loading and Preprocessing ---
def load_and_prepare_data():
    """Loads and preprocesses the data."""
    global df_global, unique_waves, unique_tas, unique_countries
    if not os.path.exists(DATA_FILE):
        logging.error(f"Data file not found: {DATA_FILE}")
        raise FileNotFoundError(f"Data file not found: {DATA_FILE}")
    try:
        df_global = pd.read_csv(DATA_FILE)
        logging.info(f"Loaded data from {DATA_FILE}. Shape: {df_global.shape}")

        # Clean Country and add Country_ID
        if 'Country' in df_global.columns:
            df_global['Country'] = df_global['Country'].replace('-', 'G5').fillna('Unknown')
            unique_countries = sorted(df_global['Country'].astype(str).unique().tolist())
            df_global['Country_ID'] = df_global['Country'].apply(map_country_to_id)
            logging.info("Cleaned 'Country', created 'Country_ID'. Countries: %s, IDs: %s", len(unique_countries), df_global['Country_ID'].nunique())
        else:
             logging.warning("'Country' column missing. Cannot create Country_ID.")
             df_global['Country_ID'] = 'N/A'

        # Parse Wave Date
        if 'Wave Date' in df_global.columns:
             df_global['Wave Date'] = pd.to_datetime(df_global['Wave Date'], errors='coerce')
             invalid_dates = df_global['Wave Date'].isna().sum()
             if invalid_dates > 0: logging.warning(f"{invalid_dates} 'Wave Date' entries failed parsing.")
             logging.info("Parsed 'Wave Date' column.")
        else: logging.warning("'Wave Date' column missing.")

        # Clean Net Score
        if 'Net Score' in df_global.columns:
            df_global['Net Score'] = df_global['Net Score'].apply(clean_net_score)
            logging.info("Cleaned 'Net Score' column.")
        else: logging.warning("'Net Score' column missing.")

        # Extract unique filters
        if 'Wave' in df_global.columns: unique_waves = sorted(df_global['Wave'].dropna().astype(str).unique().tolist())
        if 'TA' in df_global.columns: unique_tas = sorted(df_global['TA'].dropna().astype(str).unique().tolist())

    except Exception as e:
        logging.error(f"Error loading/processing data: {e}", exc_info=True)
        df_global = None

# --- Flask Routes ---
@app.route('/')
def index():
    """Displays the main filter form."""
    if df_global is None:
        return render_template('error.html', message="Error: Data could not be loaded.")
    return render_template('index.html', waves=unique_waves, tas=unique_tas, countries=unique_countries)

@app.route('/get_filters_for_ta', methods=['POST'])
def get_filters_for_ta():
    """Gets unique Brands and Companies for selected TAs."""
    # (Code for this route remains the same as your working version)
    if df_global is None: return jsonify({"error": "Data not loaded on server"}), 500
    selected_tas = request.json.get('tas', [])
    if not selected_tas: return jsonify({"brands": [], "companies": []})
    try:
        df_filtered_ta = df_global[df_global['TA'].isin(selected_tas)]
        brands = sorted([b for b in df_filtered_ta['Brand'].dropna().astype(str).unique() if b])
        companies = sorted([c for c in df_filtered_ta['Company'].dropna().astype(str).unique() if c])
        return jsonify({"brands": brands, "companies": companies})
    except Exception as e:
        logging.error(f"Error getting filters for TA: {e}", exc_info=True)
        return jsonify({"error": f"Internal server error: {str(e)}"}), 500


# ==============================================================================
# --- Main Route for Generation (Complete, Integrated Version) ---
# ==============================================================================

@app.route('/generate', methods=['POST'])
def generate_presentation():
    """
    Opens template, finds slides/charts, updates MAC/NPS/Equity based on current_wave,
    updates MVP charts (with Month-Year axis) based on selected mvp_trend_waves,
    then applies Excel overrides.
    """
    # --- Initial Checks ---
    if df_global is None:
        return render_template('error.html', message="Error: Source data not available."), 500
    if not os.path.exists(TEMPLATE_FILE):
        return render_template('error.html', message=f"Error: Template file '{TEMPLATE_FILE}' not found."), 500
    # Check for essential columns needed later
    required_cols = ['TA', 'Wave', 'Country', 'Metric', 'Net Score', 'Company', 'Brand', 'Indication', 'Wave Date', 'Country_ID']
    missing_cols = [col for col in required_cols if col not in df_global.columns]
    if missing_cols:
        logging.error(f"Essential columns missing post-load: {missing_cols}")
        return render_template('error.html', message=f"Error: Missing required columns in source data: {', '.join(missing_cols)}."), 500

    try:
        # --- Get Form Data ---
        current_wave = request.form.get('current_wave')
        prior_wave = request.form.get('prior_wave') # Optional reference
        mvp_trend_waves = request.form.getlist('mvp_trend_waves')
        ta_selections = parse_dynamic_form_data(request.form)
        selected_countries = request.form.getlist('countries')

        # --- Validation ---
        if not current_wave:
             return render_template('error.html', message="Error: Please select a Current Wave."), 400
        if not mvp_trend_waves:
             return render_template('error.html', message="Error: Please select at least one Wave for the MVP Trend."), 400
        if not ta_selections:
             return render_template('error.html', message="Error: Please select at least one TA and Brands/Companies."), 400

        if not selected_countries:
             selected_countries = unique_countries # Use all if none selected
             logging.info("Using all available countries as none were selected.")
        else:
             logging.info(f"Selected countries: {selected_countries}")
        logging.info(f"Selected Current Wave (Standard Charts): {current_wave}")
        logging.info(f"Selected MVP Trend Waves: {mvp_trend_waves}")

        # --- Load the EXISTING presentation ---
        prs = Presentation(TEMPLATE_FILE)
        logging.info(f"Loaded template presentation: '{TEMPLATE_FILE}'")

        # --- Process Each Selected TA ---
        tas_processed_count = 0
        for ta, selections in ta_selections.items():
            logging.info(f"--- Processing TA: {ta} ---")
            selected_brands = selections.get('brands', [])
            selected_companies = selections.get('companies', [])

            if not selected_brands and not selected_companies:
                 logging.warning(f"Skipping TA '{ta}' as no brands or companies were selected for it.")
                 continue

            # --- Find the slide for this TA ---
            target_title_shape_name = f"{SLIDE_TITLE_SHAPE_PREFIX}{ta}"
            target_slide = find_slide_by_title_shape(prs, target_title_shape_name)
            if target_slide is None:
                 logging.error(f"Slide for TA '{ta}' not found using shape '{target_title_shape_name}'. Skipping this TA.")
                 continue # Skip this TA if slide isn't found

            logging.info(f"Found target slide for TA '{ta}'.")
            tas_processed_count += 1 # Increment only if slide found

            # --- Update TA Name Textbox ---
            try:
                ta_name_textbox = find_shape_on_slide(target_slide, TA_NAME_TEXTBOX_NAME)
                if ta_name_textbox and hasattr(ta_name_textbox, 'text_frame'):
                    ta_name_textbox.text_frame.text = ta
                    logging.info(f"Updated '{TA_NAME_TEXTBOX_NAME}' textbox for {ta}.")
                elif ta_name_textbox:
                    logging.warning(f"Shape '{TA_NAME_TEXTBOX_NAME}' found for {ta} but has no text_frame.")
                # else: find_shape_on_slide logs warning
            except Exception as e:
                 logging.error(f"Error updating '{TA_NAME_TEXTBOX_NAME}' for {ta}: {e}")

            # --- Prepare & Update STANDARD Charts (MAC, NPS, EQUITY using current_wave) ---
            logging.info(f"--- Updating Standard Charts for TA: {ta} using Wave: {current_wave} ---")
            try: # Wrap standard chart processing for this TA
                df_ta_current_wave_countries = df_global[
                    (df_global['TA'] == ta) &
                    (df_global['Wave'] == current_wave) &
                    (df_global['Country'].isin(selected_countries))
                ].copy()

                if df_ta_current_wave_countries.empty:
                     logging.warning(f"No data for TA '{ta}', Wave '{current_wave}'. Skipping standard charts.")
                else:
                    # MAC Chart
                    pivot_mac = pd.DataFrame()
                    if selected_companies:
                        df_mac = df_ta_current_wave_countries[(df_ta_current_wave_countries['Metric'] == 'MOST ACTIVE COMMUNICATOR') & (df_ta_current_wave_countries['Company'].isin(selected_companies))].copy()
                        df_mac['Net Score'] = pd.to_numeric(df_mac['Net Score'], errors='coerce')
                        if not df_mac.empty and not df_mac['Net Score'].isna().all():
                             try: pivot_mac = pd.pivot_table(df_mac, index='Country', columns='Company', values='Net Score', aggfunc='first').fillna(0).reset_index()
                             except Exception as e: logging.error(f"Error creating MAC pivot for {ta}: {e}")
                    mac_chart_name = f"{ta}{MAC_CHART_SUFFIX}"
                    mac_chart_shape = find_shape_on_slide(target_slide, mac_chart_name)
                    if mac_chart_shape and mac_chart_shape.has_chart:
                         if not pivot_mac.empty:
                              chart_data_mac = prepare_chart_data(pivot_mac)
                              if chart_data_mac.categories:
                                 try: mac_chart_shape.chart.replace_data(chart_data_mac);logging.info(f"Updated chart '{mac_chart_name}'.")
                                 except Exception as e_rep: logging.error(f"Error replacing data for {mac_chart_name}: {e_rep}")
                                        # --- Add specific formatting for J&J series ---
                                 try:
                                                chart = mac_chart_shape.chart
                                                target_series_name = "J&J" # The exact name of the series
                                                found_series = False
                                                for series in chart.series:
                                                    # Compare series name, handle potential NoneType if name is missing
                                                    if hasattr(series, 'name') and series.name == target_series_name:
                                                        logging.info(f"Applying custom formatting to '{target_series_name}' series in chart '{mac_chart_name}'.")

                                                        # --- Line Formatting ---
                                                        #line = series.format.line
                                                        # Set line color to Red
                                                        #line.color.rgb = RGBColor(255, 0, 0)
                                                        # Optional: Adjust line width if desired
                                                        # line.width = Pt(2.0)

                                                        # --- Marker Formatting ---
                                                        marker = series.marker
                                                        # Set marker style to Diamond
                                                        marker.style = XL_MARKER_STYLE.DIAMOND
                                                        # Set marker size (adjust as needed, default is often 5)
                                                        marker.size = 7 # Example size

                                                        # Optional: Set marker fill and border color to match the line
                                                        # Solid fill for the marker
                                                        marker.format.fill.solid()
                                                        marker.format.fill.fore_color.rgb = RGBColor(255, 0, 0)
                                                        # Marker border color
                                                        marker.format.line.color.rgb = RGBColor(255, 0, 0)
                                                        # Optional: Marker border width
                                                        # marker.format.line.width = Pt(1.0)

                                                        found_series = True
                                                        # Optional: break if you only expect one J&J series
                                                        break

                                                if not found_series:
                                                    logging.warning(f"Series '{target_series_name}' not found in chart '{mac_chart_name}' after data update. Cannot apply custom formatting.")

                                 except AttributeError as ae:
                                            # Catch errors if format attributes don't exist (less likely for standard charts)
                                            logging.error(f"Attribute error applying formatting to '{target_series_name}' in chart '{mac_chart_name}': {ae}. Chart might not support expected formatting.", exc_info=True)
                                 except Exception as format_err:
                                            logging.error(f"General error applying custom formatting to '{target_series_name}' in chart '{mac_chart_name}': {format_err}", exc_info=True)
                                            # --- End J&J formatting ---


                                
                              else: logging.warning(f"Prepared data empty for '{mac_chart_name}'.")
                         else: logging.info(f"Pivot empty for '{mac_chart_name}'. Not updated.")
                    # else: find_shape logs warning

                    # NPS Chart
                    pivot_nps = pd.DataFrame()
                    if selected_brands:
                        df_nps = df_ta_current_wave_countries[(df_ta_current_wave_countries['Metric'] == 'BRAND NPS') & (df_ta_current_wave_countries['Indication'] == 'NET ABOVE LINE') & (df_ta_current_wave_countries['Brand'].isin(selected_brands))].copy()
                        df_nps['Net Score'] = pd.to_numeric(df_nps['Net Score'], errors='coerce')
                        if not df_nps.empty and not df_nps['Net Score'].isna().all():
                             try: pivot_nps = pd.pivot_table(df_nps, index='Country', columns='Brand', values='Net Score', aggfunc='first').fillna(0).reset_index()
                             except Exception as e: logging.error(f"Error creating NPS pivot for {ta}: {e}")
                    nps_chart_name = f"{ta}{NPS_CHART_SUFFIX}"
                    nps_chart_shape = find_shape_on_slide(target_slide, nps_chart_name)
                    if nps_chart_shape and nps_chart_shape.has_chart:
                         if not pivot_nps.empty:
                              chart_data_nps = prepare_chart_data(pivot_nps)
                              if chart_data_nps.categories:
                                   try: nps_chart_shape.chart.replace_data(chart_data_nps); logging.info(f"Updated chart '{nps_chart_name}'.")
                                   except Exception as e_rep: logging.error(f"Error replacing data for {nps_chart_name}: {e_rep}")
                              else: logging.warning(f"Prepared data empty for '{nps_chart_name}'.")
                         else: logging.info(f"Pivot empty for '{nps_chart_name}'. Not updated.")
                    # else: find_shape logs warning

                    # EQUITY Chart
                    pivot_equity = pd.DataFrame()
                    if selected_brands:
                        df_equity = df_ta_current_wave_countries[(df_ta_current_wave_countries['Metric'] == 'BRAND EQUITY') & (df_ta_current_wave_countries['Indication'] == 'NET ABOVE LINE') & (df_ta_current_wave_countries['Brand'].isin(selected_brands))].copy()
                        df_equity['Net Score'] = pd.to_numeric(df_equity['Net Score'], errors='coerce')
                        if not df_equity.empty and not df_equity['Net Score'].isna().all():
                             try: pivot_equity = pd.pivot_table(df_equity, index='Country', columns='Brand', values='Net Score', aggfunc='first').fillna(0).reset_index()
                             except Exception as e: logging.error(f"Error creating Equity pivot for {ta}: {e}")
                    equity_chart_name = f"{ta}{EQUITY_CHART_SUFFIX}"
                    equity_chart_shape = find_shape_on_slide(target_slide, equity_chart_name)
                    if equity_chart_shape and equity_chart_shape.has_chart:
                         if not pivot_equity.empty:
                              chart_data_equity = prepare_chart_data(pivot_equity)
                              if chart_data_equity.categories:
                                   try: equity_chart_shape.chart.replace_data(chart_data_equity); logging.info(f"Updated chart '{equity_chart_name}'.")
                                   except Exception as e_rep: logging.error(f"Error replacing data for {equity_chart_name}: {e_rep}")
                              else: logging.warning(f"Prepared data empty for '{equity_chart_name}'.")
                         else: logging.info(f"Pivot empty for '{equity_chart_name}'. Not updated.")
                    # else: find_shape logs warning

            except Exception as std_chart_err:
                logging.error(f"Error during standard chart processing for TA {ta}: {std_chart_err}", exc_info=True)


            # --- Prepare & Update NEW MVP Charts (using mvp_trend_waves) ---
            logging.info(f"--- Starting MVP Chart Updates for TA: {ta} using Waves: {mvp_trend_waves} ---")
            try: # Wrap MVP chart section for this TA
                df_ta_mvp_waves_countries = df_global[
                    (df_global['TA'] == ta) &
                    (df_global['Wave'].isin(mvp_trend_waves)) &
                    (df_global['Country'].isin(selected_countries))
                ].copy()

                if df_ta_mvp_waves_countries.empty:
                     logging.warning(f"No base data for MVP trend (TA '{ta}', Waves '{mvp_trend_waves}'). Skipping MVP updates.")
                elif not selected_companies:
                     logging.warning(f"No companies selected for TA {ta}. Skipping MVP updates.")
                else:
                    unique_country_ids_in_data = df_ta_mvp_waves_countries['Country_ID'].unique()
                    logging.info(f"Found Country IDs for MVP charts in TA {ta}: {list(unique_country_ids_in_data)}")

                    for country_id in unique_country_ids_in_data:
                        logging.debug(f"Processing MVP chart for Country ID: {country_id}")
                        try: # Wrap individual country MVP processing
                            # Filter further for MVP metric, current Country ID, and selected companies
                            df_mvp_country = df_ta_mvp_waves_countries[
                                (df_ta_mvp_waves_countries['Metric'] == 'MOST VALUABLE') &
                                (df_ta_mvp_waves_countries['Company'].isin(selected_companies)) &
                                (df_ta_mvp_waves_countries['Country_ID'] == country_id) &
                                (df_ta_mvp_waves_countries['Wave Date'].notna()) # Requires Wave Date
                            ].copy()

                            # --- Debug Log for df_mvp_country (Input to Pivot) ---
                            if not df_mvp_country.empty:
                                logging.info(f"--- Input data (df_mvp_country) for {ta}_MVP_{country_id} pivot ---")
                                try:
                                    df_details = df_mvp_country[['Wave', 'Wave Date', 'Company', 'Net Score']].sort_values('Wave Date').to_string()
                                    logging.info(f"\n{df_details}\n")
                                except Exception as log_e: logging.error(f"Error logging df_mvp_country details: {log_e}")
                                logging.info(f"--- End input data (df_mvp_country) ---")
                            # --- End Debug Log ---

                            if df_mvp_country.empty:
                                logging.warning(f"No 'MOST VALUABLE' data for {ta}, Co '{selected_companies}', CtrID '{country_id}'. Skipping.")
                                continue # Skip to next country_id

                            # --- Prepare Month-Year labels and Wave-to-Month Mapping ---
                            df_mvp_country = df_mvp_country.sort_values(by='Wave Date')
                            df_mvp_country['Wave_Month_Year'] = df_mvp_country['Wave Date'].dt.strftime('%b %y')
                            wave_to_month_map = df_mvp_country.drop_duplicates(subset=['Wave']).set_index('Wave')['Wave_Month_Year']
                            logging.debug(f"Wave to Month map for {country_id}: {wave_to_month_map.to_dict()}")
                            # --- End Mapping ---

                            # --- Create MVP pivot using 'Wave' as index ---
                            pivot_mvp = pd.DataFrame()
                            try:
                                pivot_mvp = pd.pivot_table(
                                    df_mvp_country, index='Wave', columns='Company',
                                    values='Net Score', aggfunc='first'
                                ).reset_index()

                                # Sort Pivot by Wave Number (BEFORE replacing 'Wave' column)
                                try:
                                    pivot_mvp['Wave_Num'] = pivot_mvp['Wave'].str.extract(r'(\d+)').astype(int)
                                    pivot_mvp = pivot_mvp.sort_values('Wave_Num').drop(columns='Wave_Num')
                                except Exception as sort_e:
                                    pivot_mvp = pivot_mvp.sort_values('Wave') # Fallback sort
                                    logging.warning(f"Could not sort MVP pivot numerically for {country_id}, used default sort: {sort_e}")

                                # Replace 'Wave' column with 'Month-Year' labels
                                if 'Wave' in pivot_mvp.columns:
                                    pivot_mvp['Wave'] = pivot_mvp['Wave'].map(wave_to_month_map)
                                    pivot_mvp.rename(columns={'Wave': 'Month-Year'}, inplace=True) # Rename for clarity
                                    logging.info(f"Replaced Wave with Month-Year labels in pivot_mvp for {ta}_MVP_{country_id}")
                                else:
                                    logging.warning("Could not find 'Wave' column in pivot_mvp to replace.")

                                logging.info(f"MVP Pivot Table processed for TA {ta}, Country_ID {country_id}")

                                # --- Debug Log for pivot_mvp (Final Output for Chart) ---
                                if not pivot_mvp.empty:
                                    logging.info(f"--- Contents of final pivot_mvp for {ta}_MVP_{country_id} chart ---")
                                    try:
                                        pivot_details = pivot_mvp.to_string()
                                        logging.info(f"\n{pivot_details}\n")
                                    except Exception as log_e: logging.error(f"Error converting pivot_mvp to string: {log_e}")
                                    logging.info(f"--- End contents of final pivot_mvp ---")
                                # --- End Debug Log ---

                            except Exception as e_pivot:
                                logging.error(f"Error creating/processing MVP pivot for {ta}, CtrID {country_id}: {e_pivot}", exc_info=True)
                                continue # Skip to next country ID if pivot fails

                            # --- Find and update the MVP chart ---
                            mvp_chart_name = f"{ta}{MVP_CHART_IDENTIFIER}{country_id}"
                            mvp_chart_shape = find_shape_on_slide(target_slide, mvp_chart_name)

                            if mvp_chart_shape and mvp_chart_shape.has_chart:
                                 if not pivot_mvp.empty:
                                      chart_data_mvp = prepare_chart_data(pivot_mvp) # Uses first column ('Month-Year')
                                      if chart_data_mvp.categories:
                                           try:
                                                mvp_chart_shape.chart.replace_data(chart_data_mvp)
                                                logging.info(f"Successfully updated chart '{mvp_chart_name}'.")
                                           except Exception as e_replace:
                                                logging.error(f"Error replacing data for chart '{mvp_chart_name}': {e_replace}", exc_info=True)
                                      else:
                                           logging.warning(f"Prepared chart data for '{mvp_chart_name}' has no categories. Update skipped.")
                                 else:
                                      logging.info(f"Pivot data empty for '{mvp_chart_name}'. Chart not updated.")
                                      # Consider clearing chart here if needed
                            # else: find_shape_on_slide logs warning

                        except Exception as country_err:
                            logging.error(f"Error processing MVP chart for {ta}, Country ID {country_id}: {country_err}", exc_info=True)
                            # Continue to next country ID

            except Exception as mvp_chart_err:
                 logging.error(f"General error during MVP chart processing for TA {ta}: {mvp_chart_err}", exc_info=True)
                 # Continue processing TAs and overrides

        # --- End of TA loop ---


        # --- Excel Override Section ---
        logging.info("--- Starting Excel Override Process ---")
        override_success = 0 # Initialize here
        try:
            override_file_path = OVERRIDE_EXCEL_FILE
            df_control = read_override_control_sheet(override_file_path)

            if df_control is not None:
                override_attempts = 0
                # override_success = 0 # Moved initialization before try block
                for index, row in df_control.iterrows():
                    try: # Wrap processing for single row
                        update_flag_val = row.get('Update_Flag')
                        chart_name = row.get('Chart_Name')

                        if should_update(update_flag_val) and chart_name and isinstance(chart_name, str):
                            override_attempts += 1
                            logging.info(f"Override requested for chart: '{chart_name}'")

                            df_override_data = read_named_range_data(override_file_path, chart_name)
                            if df_override_data is None or df_override_data.empty:
                                logging.warning(f"Skipping override for '{chart_name}': Data empty/not found in Excel.")
                                continue

                            chart_shape = find_shape_anywhere(prs, chart_name)
                            if not chart_shape or not chart_shape.has_chart:
                                logging.warning(f"Skipping override for '{chart_name}': Chart shape not found/not a chart in PPT.")
                                continue

                            chart_data_override = prepare_chart_data(df_override_data)
                            if chart_data_override.categories:
                                try:
                                    chart_shape.chart.replace_data(chart_data_override)
                                    logging.info(f"Successfully updated chart '{chart_name}' from Excel override.")
                                    override_success += 1
                                except Exception as e_ovr_rep:
                                    logging.error(f"Error replacing data for overridden chart '{chart_name}': {e_ovr_rep}")
                            else:
                                logging.warning(f"Prepared override data empty for '{chart_name}'. Update skipped.")
                        elif update_flag_val is not None and (not chart_name or not isinstance(chart_name, str)):
                             logging.warning(f"Row {index+2} in control sheet has update flag but invalid/missing chart name.")
                        # else: # No update flag set, do nothing
                        #     pass

                    except Exception as row_err:
                        chart_name_for_log = chart_name if 'chart_name' in locals() and chart_name else f"row {index+2}"
                        logging.error(f"Error processing override for {chart_name_for_log}: {row_err}", exc_info=True)
                        # Continue to next row

                logging.info(f"Excel Override Process Summary: Attempted={override_attempts}, Succeeded={override_success}")
            else:
                 logging.info("Skipping Excel override process (control file not found or invalid).")
        except Exception as excel_override_err:
             logging.error(f"Major error during Excel Override processing: {excel_override_err}", exc_info=True)


        # --- Final Checks & Save ---
        if tas_processed_count == 0:
            logging.warning("Warning: No TA slides were found or processed based on selections.")
            # It might be okay to return the (potentially overridden) original template
            # Or return an error if absolutely no changes were expected/made:
            # return render_template('error.html', message="Warning: No TA slides could be processed."), 404

        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        logging.info(f"Presentation update process completed. {tas_processed_count} TA slides processed via script logic.")

        return send_file(
            file_stream,
            as_attachment=True,
            download_name=OUTPUT_FILENAME,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    # --- Exception Handling ---
    except FileNotFoundError as e:
        logging.error(f"File Not Found Error: {e}")
        return render_template('error.html', message=str(e)), 500
    except KeyError as e:
         logging.error(f"KeyError during processing: {e}", exc_info=True)
         return render_template('error.html', message=f"Error: Missing expected data field: {e}"), 400
    except Exception as e:
        # Log the full traceback for unexpected errors
        logging.exception("An unexpected error occurred during presentation generation.")
        return render_template('error.html', message=f"An unexpected server error occurred: {str(e)}"), 500

# --- Main Execution ---
if __name__ == '__main__':
    load_and_prepare_data()
    if df_global is not None:
        # Use Waitress for production:
        # from waitress import serve
        # serve(app, host='0.0.0.0', port=8080)
        app.run(debug=True, port=5002) # Adjust port if needed
    else:
        print("CRITICAL: Failed to load data. Flask application cannot start.")